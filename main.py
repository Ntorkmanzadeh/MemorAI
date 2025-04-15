from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Depends, Body
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
import json
from typing import Optional, List, Dict, Any
import PyPDF2
from pptx import Presentation
import io
import logging
import os
from datetime import datetime
from pathlib import Path
import database as db
import modal
import re
import asyncio
from PIL import Image, ImageDraw
from pdf2image import convert_from_bytes

# Create logs directory if it doesn't exist
logs_dir = Path("logs")
logs_dir.mkdir(exist_ok=True)

# Set up logging
log_file = logs_dir / f"flashcard_generator_{datetime.now().strftime('%Y%m%d_%H%M%S')}.log"
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler(log_file),
        logging.StreamHandler()  # Also print to console
    ]
)
logger = logging.getLogger(__name__)

app = FastAPI()

# Configure CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Get reference to the deployed Modal function
run_ollama_prompt = modal.Function.from_name("flashcard-generator", "run_ollama_prompt")
process_image_with_llama = modal.Function.from_name("flashcard-generator", "process_image_with_llama")
process_multiple_images_with_llama = modal.Function.from_name("flashcard-generator", "process_multiple_images_with_llama")

# User dependency
async def get_current_user(user_id: str = Form(...)):
    """Get the current user from the user_id form field."""
    try:
        user = db.get_user(user_id)
        if not user:
            raise HTTPException(status_code=404, detail="User not found")
        return {"id": user_id, "name": user["name"]}
    except Exception as e:
        logger.error(f"Error getting user: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting user: {str(e)}")

def extract_text_from_pdf(file_content: bytes) -> str:
    logger.info("Extracting text from PDF")
    pdf_file = io.BytesIO(file_content)
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text() + "\n"
    logger.info(f"Extracted {len(text)} characters from PDF")
    return text

def extract_text_from_pptx(file_content: bytes) -> str:
    logger.info("Extracting text from PPTX")
    pptx_file = io.BytesIO(file_content)
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    logger.info(f"Extracted {len(text)} characters from PPTX")
    return text

def convert_pptx_to_images(pptx_content: bytes) -> List[Image.Image]:
    """Convert PPTX content to a list of PIL Images."""
    logger.info("Converting PPTX to images")
    pptx_file = io.BytesIO(pptx_content)
    prs = Presentation(pptx_file)
    images = []
    
    for slide in prs.slides:
        # Create a new image for each slide
        img = Image.new('RGB', (1920, 1080), color='white')
        draw = ImageDraw.Draw(img)
        
        # Add text from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                # Calculate text position (simplified)
                text = shape.text
                text_width = draw.textlength(text)
                text_position = ((1920 - text_width) / 2, 540)  # Center text
                draw.text(text_position, text, fill='black')
        
        images.append(img)
    
    logger.info(f"Converted {len(images)} slides to images")
    return images

def convert_pdf_to_images(pdf_content: bytes) -> List[Image.Image]:
    """Convert PDF content to a list of PIL Images."""
    logger.info("Converting PDF to images")
    return convert_from_bytes(pdf_content)

async def generate_flashcards(text: str, num_cards: int = 20) -> List[Dict[str, str]]:
    """Generate flashcards from text using Ollama."""
    logging.info(f"Starting flashcard generation for text of length {len(text)}")
    logging.info(f"Requesting {num_cards} flashcards")
    
    prompt = f"""Generate {num_cards} high-quality flashcards from the following text. Focus on key concepts, definitions, and facts that are likely to appear on a test or exam.

Guidelines for the flashcards:
1. Focus on important concepts, definitions, and key facts
2. Include specific details, numbers, and dates that might be tested
3. Create questions that test understanding rather than just memorization
4. Include both basic recall questions and more complex application questions
5. Make sure answers are concise but complete
6. Cover a range of difficulty levels
7. Include questions about relationships between concepts
8. Focus on information that would be most valuable in a test setting

CRITICAL INSTRUCTION: You must return ONLY a valid JSON array containing objects with EXACTLY these two fields: "question" and "answer".
DO NOT include any other fields like "title", "creative_title", "reference_list", or "in_text_citations".
DO NOT include any additional text, markdown formatting, or explanations.

Example format:
[
    {{"question": "What is the key concept of X and how does it relate to Y?", "answer": "X is a fundamental principle that... It relates to Y by..."}},
    {{"question": "What are the three main components of Z and their functions?", "answer": "The three components are: 1) A - responsible for..., 2) B - handles..., 3) C - manages..."}}
]

Text to generate flashcards from:
{text}

Remember: Return ONLY the JSON array with objects containing ONLY "question" and "answer" fields. You also must generate at least {num_cards} cards."""

    try:
        logging.info("Calling Modal function run_ollama_prompt...")
        # Call the Modal function directly
        response = await run_ollama_prompt.remote.aio(prompt)
        logging.info("Received response from Modal function")
        logging.info("Complete model response:")
        logging.info(response)
        
        if not response:
            logging.error("Received empty response from Modal app")
            return generate_fallback_cards(num_cards)
        
        # Try to find a JSON array in the response
        try:
            # First try to parse the entire response as JSON
            try:
                cards = json.loads(response)
                if isinstance(cards, list):
                    # Check if the cards have the correct format
                    valid_cards = []
                    for card in cards:
                        if isinstance(card, dict) and "question" in card and "answer" in card:
                            valid_cards.append({
                                "question": card["question"],
                                "answer": card["answer"]
                            })
                    
                    if valid_cards:
                        logging.info(f"Successfully parsed JSON array with {len(valid_cards)} valid cards")
                        return valid_cards
                    else:
                        logging.warning("Parsed JSON array but no valid cards found")
            except json.JSONDecodeError:
                # If that fails, try to extract JSON from the response
                logging.info("Failed to parse entire response as JSON, trying to extract JSON array")
            
            # Look for text between square brackets, handling potential markdown code blocks
            # First try to find a JSON array in a code block
            code_block_match = re.search(r'```(?:json)?\s*(\[[\s\S]*?\])\s*```', response)
            if code_block_match:
                json_str = code_block_match.group(1)
                try:
                    cards = json.loads(json_str)
                    if isinstance(cards, list):
                        # Check if the cards have the correct format
                        valid_cards = []
                        for card in cards:
                            if isinstance(card, dict) and "question" in card and "answer" in card:
                                valid_cards.append({
                                    "question": card["question"],
                                    "answer": card["answer"]
                                })
                        
                        if valid_cards:
                            logging.info(f"Successfully extracted JSON from code block with {len(valid_cards)} valid cards")
                            return valid_cards
                        else:
                            logging.warning("Extracted JSON from code block but no valid cards found")
                except json.JSONDecodeError:
                    logging.warning("Failed to parse JSON from code block")
            
            # If no code block found or parsing failed, try to find any JSON array
            match = re.search(r'(\[[\s\S]*?\])', response)
            if match:
                json_str = match.group(1)
                try:
                    cards = json.loads(json_str)
                    if isinstance(cards, list):
                        # Check if the cards have the correct format
                        valid_cards = []
                        for card in cards:
                            if isinstance(card, dict) and "question" in card and "answer" in card:
                                valid_cards.append({
                                    "question": card["question"],
                                    "answer": card["answer"]
                                })
                        
                        if valid_cards:
                            logging.info(f"Successfully extracted JSON array with {len(valid_cards)} valid cards")
                            return valid_cards
                        else:
                            logging.warning("Extracted JSON array but no valid cards found")
                except json.JSONDecodeError:
                    logging.warning("Failed to parse extracted JSON array")
            
            # If all extraction attempts fail, try to clean the response and parse it
            cleaned_response = re.sub(r'```.*?```', '', response, flags=re.DOTALL)  # Remove code blocks
            cleaned_response = re.sub(r'^.*?\[', '[', cleaned_response, flags=re.DOTALL)  # Remove text before first [
            cleaned_response = re.sub(r'\].*?$', ']', cleaned_response, flags=re.DOTALL)  # Remove text after last ]
            
            try:
                cards = json.loads(cleaned_response)
                if isinstance(cards, list):
                    # Check if the cards have the correct format
                    valid_cards = []
                    for card in cards:
                        if isinstance(card, dict) and "question" in card and "answer" in card:
                            valid_cards.append({
                                "question": card["question"],
                                "answer": card["answer"]
                            })
                    
                    if valid_cards:
                        logging.info(f"Successfully parsed cleaned JSON with {len(valid_cards)} valid cards")
                        return valid_cards
                    else:
                        logging.warning("Parsed cleaned JSON but no valid cards found")
            except json.JSONDecodeError:
                logging.warning("Failed to parse cleaned JSON")
            
            # If we still don't have valid cards, try to extract question-answer pairs from the text
            logging.info("Attempting to extract question-answer pairs from text")
            qa_pairs = re.findall(r'["\']question["\']\s*:\s*["\']([^"\']+)["\']\s*,\s*["\']answer["\']\s*:\s*["\']([^"\']+)["\']', response)
            if qa_pairs:
                valid_cards = [{"question": q, "answer": a} for q, a in qa_pairs]
                logging.info(f"Successfully extracted {len(valid_cards)} question-answer pairs from text")
                return valid_cards
            
            logging.warning("No valid JSON array or question-answer pairs found in response after all attempts")
            return generate_fallback_cards(num_cards)
            
        except Exception as e:
            logging.error(f"Error processing JSON: {str(e)}")
            return generate_fallback_cards(num_cards)
            
    except Exception as e:
        logging.error(f"Error generating flashcards: {str(e)}")
        logging.error(f"Error type: {type(e)}")
        import traceback
        logging.error(f"Traceback: {traceback.format_exc()}")
        return generate_fallback_cards(num_cards)

def generate_fallback_cards(num_cards: int) -> List[Dict[str, str]]:
    """Generate fallback flashcards when the AI service fails."""
    logging.info(f"Generating {num_cards} fallback flashcards")
    cards = []
    for i in range(num_cards):
        cards.append({
            "question": f"Sample question {i+1}?",
            "answer": f"This is a sample answer {i+1}. The actual AI service is currently unavailable."
        })
    return cards

# User management endpoints
@app.post("/users")
async def create_user(name: str = Form(...)):
    """Create a new user."""
    try:
        user_id = db.create_user(name)
        return {"user_id": user_id, "name": name}
    except Exception as e:
        logger.error(f"Error creating user: {e}")
        raise HTTPException(status_code=500, detail=f"Error creating user: {str(e)}")

@app.get("/users/{user_id}")
async def get_user(user_id: str):
    """Get a user by ID."""
    user = db.get_user(user_id)
    if not user:
        raise HTTPException(status_code=404, detail="User not found")
    return user

# Deck management endpoints
@app.post("/decks")
async def create_deck(user: Dict[str, Any] = Depends(get_current_user), name: str = Form(...), description: str = Form("")):
    """Create a new deck for a user."""
    try:
        deck_id = db.create_deck(user["id"], name, description)
        return {
            "deck_id": deck_id,
            "name": name,
            "description": description,
            "user_id": user["id"]
        }
    except Exception as e:
        logger.error(f"Error creating deck: {e}")
        raise HTTPException(status_code=500, detail=f"Error creating deck: {str(e)}")

@app.get("/decks")
async def get_decks(user: Dict[str, Any] = Depends(get_current_user)):
    """Get all decks for a user."""
    try:
        decks = db.get_decks(user["id"])
        return {"decks": decks}
    except Exception as e:
        logger.error(f"Error getting decks: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting decks: {str(e)}")

@app.put("/decks/{deck_id}")
async def update_deck(deck_id: str, name: str = Form(...), description: str = Form("")):
    """Update a deck."""
    try:
        db.update_deck(deck_id, name, description)
        return {"deck_id": deck_id, "name": name, "description": description}
    except Exception as e:
        logger.error(f"Error updating deck: {e}")
        raise HTTPException(status_code=500, detail=f"Error updating deck: {str(e)}")

@app.delete("/decks/{deck_id}")
async def delete_deck(deck_id: str):
    """Delete a deck."""
    try:
        db.delete_deck(deck_id)
        return {"message": "Deck deleted successfully"}
    except Exception as e:
        logger.error(f"Error deleting deck: {e}")
        raise HTTPException(status_code=500, detail=f"Error deleting deck: {str(e)}")

# Flashcard management endpoints
@app.get("/decks/{deck_id}/flashcards")
async def get_flashcards(deck_id: str):
    """Get all flashcards in a deck."""
    try:
        flashcards = db.get_flashcards(deck_id)
        return {"flashcards": flashcards}
    except Exception as e:
        logger.error(f"Error getting flashcards: {e}")
        raise HTTPException(status_code=500, detail=f"Error getting flashcards: {str(e)}")

@app.post("/decks/{deck_id}/flashcards")
async def create_flashcard(deck_id: str, question: str = Form(...), answer: str = Form(...)):
    """Create a new flashcard in a deck."""
    try:
        card_id = db.create_flashcard(deck_id, question, answer)
        return {"card_id": card_id, "question": question, "answer": answer}
    except Exception as e:
        logger.error(f"Error creating flashcard: {e}")
        raise HTTPException(status_code=500, detail=f"Error creating flashcard: {str(e)}")

@app.put("/flashcards/{card_id}")
async def update_flashcard(card_id: str, question: str = Form(...), answer: str = Form(...)):
    """Update a flashcard."""
    try:
        db.update_flashcard(card_id, question, answer)
        return {"card_id": card_id, "question": question, "answer": answer}
    except Exception as e:
        logger.error(f"Error updating flashcard: {e}")
        raise HTTPException(status_code=500, detail=f"Error updating flashcard: {str(e)}")

@app.delete("/flashcards/{card_id}")
async def delete_flashcard(card_id: str):
    """Delete a flashcard."""
    try:
        db.delete_flashcard(card_id)
        return {"message": "Flashcard deleted successfully"}
    except Exception as e:
        logger.error(f"Error deleting flashcard: {e}")
        raise HTTPException(status_code=500, detail=f"Error deleting flashcard: {str(e)}")

@app.post("/process-file")
async def process_file(
    user: Dict[str, Any] = Depends(get_current_user),
    deck_id: str = Form(...),
    file: Optional[UploadFile] = File(None),
    text: Optional[str] = Form(None)
):
    try:
        content_text = ""
        
        if file:
            logger.info(f"Processing file: {file.filename}")
            file_content = await file.read()
            if file.filename.endswith('.pdf'):
                logger.info("Starting PDF processing")
                images = convert_pdf_to_images(file_content)
                logger.info(f"Converted PDF to {len(images)} images")
                
                # Convert all images to bytes
                image_bytes_list = []
                for image in images:
                    buffered = io.BytesIO()
                    image.save(buffered, format="PNG")
                    image_bytes_list.append(buffered.getvalue())
                
                # Define the message for flashcard generation
                message = (
                    "Below are summaries made by llama3.2-vision for a series of pages from a PDF. Please analyze the content of these pages and provide a comprehensive, detailed description that synthesizes information across all pages. "
                    "Your summary should be thorough and specific, maintaining at least 50% of the original length. "
                    "For each key concept, include:\n"
                    "- Precise definitions and terminology\n"
                    "- Step-by-step explanations of processes\n"
                    "- Specific examples and applications\n"
                    "- Relationships between different concepts\n"
                    "- Important numerical values or measurements\n\n"
                    "Note: Some pages may show sequential steps in a process or animation. In such cases, explain the complete process flow rather than repeating similar information. "
                    "This summary will be used to generate detailed flashcards, so include all specific details, formulas, and technical terms that would be important for exam preparation."
                )
                
                # Process all images at once
                logger.info(f"Processing {len(image_bytes_list)} images in batch")
                content_text = await process_multiple_images_with_llama.remote.aio(image_bytes_list, message)
                
            elif file.filename.endswith('.pptx'):
                logger.info("Starting PPTX processing")
                images = convert_pptx_to_images(file_content)
                logger.info(f"Converted PPTX to {len(images)} images")
                
                # Convert all images to bytes
                image_bytes_list = []
                for image in images:
                    buffered = io.BytesIO()
                    image.save(buffered, format="PNG")
                    image_bytes_list.append(buffered.getvalue())
                
                # Define the message for flashcard generation
                message = (
                    "Analyze the content of these slides and generate high-quality flashcards. "
                    "Focus on key concepts, definitions, and facts that are likely to appear on a test or exam. "
                    "Each flashcard should have a 'question' and an 'answer'."
                )
                
                # Process all images at once
                logger.info(f"Processing {len(image_bytes_list)} slides in batch")
                content_text = await process_multiple_images_with_llama.remote.aio(image_bytes_list, message)
                
            elif file.filename.lower().endswith(('.png', '.jpg', '.jpeg')):
                # For single images, use the existing function
                message = (
                    "Analyze the content of this image and generate high-quality flashcards. "
                    "Focus on key concepts, definitions, and facts that are likely to appear on a test or exam. "
                    "Each flashcard should have a 'question' and an 'answer'."
                )
                
                content_text = await process_image_with_llama.remote.aio(file_content, message)
            else:
                logger.warning(f"Unsupported file format: {file.filename}")
                return JSONResponse(
                    status_code=400,
                    content={"error": "Unsupported file format. Please upload PDF, PPTX, or image files (PNG, JPG, JPEG)."}
                )
        elif text:
            logger.info("Processing text input")
            content_text = text
        else:
            logger.warning("No file or text provided")
            return JSONResponse(
                status_code=400,
                content={"error": "No file or text provided"}
            )
        flashcards = await generate_flashcards(content_text)
        logger.info(f"Generated {len(flashcards)} flashcards")
        
        db.import_flashcards(deck_id, flashcards)
        
        return {"message": f"Successfully added {len(flashcards)} flashcards to deck", "flashcards": flashcards}
    
    except Exception as e:
        logger.error(f"Error processing content: {str(e)}")
        return JSONResponse(
            status_code=500,
            content={"error": f"Error processing content: {str(e)}"}
        )

@app.get("/health")
async def health_check():
    return {"status": "healthy"} 